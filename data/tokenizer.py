# tokenizer.py
import json, csv, openpyxl
from pathlib import Path
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from nltk.corpus import stopwords
import re

MAX_TOKENS = 50
STOPWORDS = set(stopwords.words("english"))

def extract_text(file_path):
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".txt":
            return _read_txt(file_path), "text"
        elif ext == ".pdf":
            return _extract_pdf(file_path), "text"
        elif ext == ".docx":
            return _extract_docx(file_path), "text"
        elif ext == ".pptx":
            return _extract_pptx(file_path), "text+notes"
        elif ext == ".xlsx":
            return _extract_xlsx(file_path), "text"
        elif ext == ".csv":
            return _extract_csv(file_path), "text"
        elif ext == ".json":
            return _extract_json(file_path), "text"
    except Exception as e:
        print(f"[ERROR] {file_path} → {e}")
    return "", "error"

def tokenize(text):
    tokens = re.findall(r'\b[a-z]+\b', text.lower())
    return [t for t in tokens if t not in STOPWORDS][:MAX_TOKENS]

# Extractors below unchanged
def _read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def _extract_pdf(file_path):
    reader = PdfReader(file_path)
    return "\n".join([page.extract_text() or "" for page in reader.pages])

def _extract_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def _extract_pptx(file_path):
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            texts.append(notes)
    return "\n".join(texts)

def _extract_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path, read_only=True)
    return "\n".join([str(cell) for sheet in wb.worksheets for row in sheet.iter_rows(values_only=True) for cell in row if cell is not None])

def _extract_csv(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return "\n".join([" ".join(row) for row in csv.reader(f)])

def _extract_json(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return _flatten_json(json.load(f))

def _flatten_json(data):
    texts = []
    def recurse(obj):
        if isinstance(obj, dict):
            for k, v in obj.items(): texts.append(str(k)); recurse(v)
        elif isinstance(obj, list): [recurse(i) for i in obj]
        else: texts.append(str(obj))
    recurse(data)
    return " ".join(texts)
