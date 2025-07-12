import os
import csv
import mimetypes
from pathlib import Path
from tqdm import tqdm
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# -------- Content Extraction Logic --------
def extract_text_from_file(filepath):
    ext = filepath.suffix.lower()
    try:
        if ext == ".pdf":
            reader = PdfReader(str(filepath))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext == ".docx":
            doc = Document(str(filepath))
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext == ".pptx":
            prs = Presentation(str(filepath))
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif ext == ".txt" or ext == ".md":
            text = filepath.read_text(encoding='utf-8', errors='ignore')
        else:
            text = ""
    except Exception as e:
        text = f"[Error reading file: {e}]"
    return text.strip()

# -------- Utility Functions --------
def extract_tags_from_filename(name):
    return [token for token in Path(name).stem.replace("_", " ").replace("-", " ").split() if len(token) > 2]

def extract_significant_content(text, max_length=1000, min_length=100):
    text = text.strip()
    if len(text) > max_length:
        return text[:max_length] + "..."
    elif len(text) < min_length:
        return text + " " * (min_length - len(text))  # padding for low-content
    else:
        return text

def is_supported_file(filepath):
    return filepath.suffix.lower() in {".pdf", ".docx", ".pptx", ".txt", ".md"}

# -------- Main Logic --------
def extract_all_metadata(input_path, output_csv="file_metadata.csv"):
    file_data = []

    for root, dirs, files in tqdm(os.walk(input_path), desc="Scanning files"):
        for filename in files:
            filepath = Path(root) / filename
            if not is_supported_file(filepath):
                continue

            metadata = {}
            metadata["file_name"] = filepath.name
            metadata["file_path"] = str(filepath)
            metadata["file_size_kb"] = round(filepath.stat().st_size / 1024, 2)
            metadata["file_type"] = filepath.suffix.lower()
            metadata["parent_folder"] = filepath.parent.name
            metadata["tags"] = ", ".join(set(extract_tags_from_filename(filename) + extract_tags_from_filename(filepath.parent.name)))

            content = extract_text_from_file(filepath)
            metadata["content_summary"] = extract_significant_content(content)

            file_data.append(metadata)

    # Write to CSV
    keys = ["file_name", "file_path", "file_size_kb", "file_type", "parent_folder", "tags", "content_summary"]
    with open(output_csv, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=keys)
        writer.writeheader()
        writer.writerows(file_data)

    print(f"\n✅ Metadata extraction complete. Results saved to: {output_csv}")

# -------- Entry Point --------
if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Extract metadata and content summary from document files.")
    parser.add_argument("path", help="Input folder path")
    parser.add_argument("--output", default="file_metadata.csv", help="CSV output path")
    args = parser.parse_args()

    input_folder = Path(args.path).resolve()
    if not input_folder.exists():
        print(f"❌ Path does not exist: {input_folder}")
    else:
        extract_all_metadata(input_folder, args.output)
