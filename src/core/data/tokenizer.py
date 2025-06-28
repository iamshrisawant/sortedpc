# data/tokenizer.py
import os
import re
from collections import Counter
from sklearn.feature_extraction.text import ENGLISH_STOP_WORDS

from docx import Document
import pdfplumber
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup
import nltk

MAX_TOKENS = 50
nltk.download('punkt', quiet=True)

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.docx':
            doc = Document(file_path)
            return '\n'.join(p.text for p in doc.paragraphs)

        elif ext == '.pdf':
            with pdfplumber.open(file_path) as pdf:
                return '\n'.join(page.extract_text() or '' for page in pdf.pages)

        elif ext == '.pptx':
            prs = Presentation(file_path)
            return '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

        elif ext == '.xlsx':
            wb = openpyxl.load_workbook(file_path, data_only=True)
            return '\n'.join(str(cell.value) for sheet in wb for row in sheet.iter_rows() for cell in row if cell.value)

        elif ext == '.html':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                return soup.get_text(separator=' ')

        else:  # Fallback to plain text
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    except Exception:
        return ''  # Failsafe fallback


def tokenize_file(file_path):
    content = extract_text(file_path)
    file_name_tokens = re.findall(r'\b\w+\b', os.path.basename(file_path).lower())
    tokens = re.findall(r'\b\w+\b', content.lower()) + file_name_tokens

    filtered = [t for t in tokens if t not in ENGLISH_STOP_WORDS and len(t) > 1]
    token_counts = Counter(filtered)
    most_common = token_counts.most_common(MAX_TOKENS)

    if len(most_common) < MAX_TOKENS:
        extras = [t for t in token_counts if t not in dict(most_common)]
        most_common += [(t, token_counts[t]) for t in extras[:MAX_TOKENS - len(most_common)]]

    return [t for t, _ in most_common]
