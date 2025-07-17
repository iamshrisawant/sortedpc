import hashlib
import string
import re
import logging
from pathlib import Path
from typing import Union, Dict

import docx
import pdfplumber
import pandas as pd
from pptx import Presentation
from openpyxl import load_workbook

from nltk.corpus import stopwords

logging.getLogger("pdfminer").setLevel(logging.ERROR)
logger = logging.getLogger(__name__)

STOPWORDS = set(stopwords.words('english'))

def extract(file_path: Union[str, Path]) -> Dict[str, str]:
    path = Path(file_path)
    if not path.is_file():
        logger.warning(f"[Extractor] Invalid file: {file_path}")
        raise FileNotFoundError(f"Path does not point to a file: {file_path}")

    file_type = path.suffix.lower().lstrip('.')
    file_name = path.stem
    parent_folder_path = path.parent.resolve()
    parent_folder_name = path.parent.name

    logger.debug(f"[Extractor] Extracting content from: {path.name}")
    content = extract_content(path, file_type)
    cleaned_content = clean_text(content)
    cleaned_file_name = clean_text(file_name)
    cleaned_folder_name = clean_text(parent_folder_name)
    content_hash = compute_hash(content)

    logger.info(f"[Extractor] {path.name} | Type: {file_type} | Hash: {content_hash[:8]} | Length: {len(cleaned_content)} chars")

    return {
        "file_path": str(path),
        "content": cleaned_content,
        "file_name": cleaned_file_name,
        "parent_folder": cleaned_folder_name,
        "parent_folder_path": str(parent_folder_path),  # ← NEW and important
        "file_type": file_type,
        "content_hash": content_hash,
    }

def extract_content(path: Path, file_type: str) -> str:
    try:
        if file_type == "pdf":
            with pdfplumber.open(path) as pdf:
                return "\n".join(page.extract_text() or "" for page in pdf.pages)

        elif file_type == "docx":
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs)

        elif file_type == "pptx":
            prs = Presentation(path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

        elif file_type == "xlsx":
            wb = load_workbook(path, data_only=True)
            return "\n".join(
                str(cell.value)
                for ws in wb.worksheets
                for row in ws.iter_rows()
                for cell in row
                if cell.value is not None
            )

        elif file_type == "csv":
            df = pd.read_csv(path)
            return df.astype(str).to_string(index=False)

        elif file_type in ("txt", "md"):
            return path.read_text(encoding="utf-8", errors="ignore")

        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    except Exception as e:
        logger.error(f"[Extractor] Failed to read {path.name} ({file_type}): {e}")
        raise RuntimeError(f"Failed to extract content from {path.name}: {e}")

def clean_text(text: str) -> str:
    text = text.lower()
    text = re.sub(r"[^\w\s]", " ", text)
    words = text.split()
    filtered = [w for w in words if w not in STOPWORDS]
    return " ".join(filtered)

def compute_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8")).hexdigest()
