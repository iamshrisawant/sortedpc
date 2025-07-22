# processor.py

import hashlib
import logging
import re
from pathlib import Path
from typing import Dict, List, Union

import docx
import numpy as np
import pandas as pd
import pdfplumber
from openpyxl import load_workbook
from pptx import Presentation
from sentence_transformers import SentenceTransformer

from nltk.corpus import stopwords

# --- Logger Setup ---
logging.getLogger("pdfminer").setLevel(logging.ERROR)
logger = logging.getLogger(__name__)

# --- Globals & Constants ---
STOPWORDS = set(stopwords.words('english'))
MODEL_NAME = "all-MiniLM-L6-v2"
_model = None  # Initialized lazily

# --------------------------------------------------------------------------
# --- TEXT EXTRACTION LOGIC (from extractor.py)
# --------------------------------------------------------------------------

def _extract_content(path: Path, file_type: str) -> str:
    """Extracts raw text content from a file."""
    try:
        if file_type == "pdf":
            with pdfplumber.open(path) as pdf:
                return "\n".join(page.extract_text() or "" for page in pdf.pages)
        elif file_type == "docx":
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        elif file_type == "pptx":
            prs = Presentation(path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif file_type == "xlsx":
            wb = load_workbook(path, data_only=True)
            return "\n".join(
                str(cell.value)
                for ws in wb.worksheets
                for row in ws.iter_rows()
                for cell in row
                if cell.value is not None
            )
        elif file_type == "csv":
            df = pd.read_csv(path)
            return df.astype(str).to_string(index=False)
        elif file_type in ("txt", "md"):
            return path.read_text(encoding="utf-8", errors="ignore")
        else:
            logger.warning(f"[Processor] Unsupported file type for extraction: {file_type}")
            return ""
    except Exception as e:
        logger.error(f"[Processor] Failed to read {path.name} ({file_type}): {e}")
        raise RuntimeError(f"Failed to extract content from {path.name}: {e}")

def _clean_text(text: str) -> str:
    """Cleans text by lowercasing, removing punctuation, and filtering stopwords."""
    text = text.lower()
    text = re.sub(r"[^\w\s]", " ", text)
    words = text.split()
    filtered = [w for w in words if w not in STOPWORDS]
    return " ".join(filtered)

def _compute_hash(text: str) -> str:
    """Computes a SHA256 hash of the text content."""
    return hashlib.sha256(text.encode("utf-8")).hexdigest()

# --------------------------------------------------------------------------
# --- TEXT CHUNKING LOGIC (from chunker.py)
# --------------------------------------------------------------------------

def _chunk_text(text: str, max_chunk_size: int = 300, overlap: int = 50) -> List[str]:
    """Splits cleaned text into overlapping word-based chunks."""
    if not text.strip():
        logger.warning("[Processor] Empty text received, returning 0 chunks.")
        return []
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + max_chunk_size, len(words))
        chunk = words[start:end]
        chunks.append(" ".join(chunk))
        start += max_chunk_size - overlap
    logger.info(f"[Processor] Chunked text into {len(chunks)} chunk(s).")
    return chunks

# --------------------------------------------------------------------------
# --- EMBEDDING LOGIC (from embedder.py)
# --------------------------------------------------------------------------

def _load_model(local_only: bool = False):
    """Lazily loads the SentenceTransformer model."""
    global _model
    if _model is not None:
        return _model
    try:
        logger.info(f"[Processor] Loading model: {MODEL_NAME} (local_only={local_only})")
        _model = SentenceTransformer(MODEL_NAME, local_files_only=local_only)
        logger.info("[Processor] Model loaded successfully.")
    except Exception as e:
        logger.error(f"[Processor] Failed to load model: {e}")
        raise
    return _model

def _embed_texts(chunks: List[str]) -> List[List[float]]:
    """Generates sentence embeddings for a list of text chunks."""
    if not chunks:
        logger.warning("[Processor] No chunks provided for embedding.")
        return []
    logger.info(f"[Processor] Generating embeddings for {len(chunks)} chunk(s)...")
    try:
        model = _load_model(local_only=True)  # Assume model is pre-downloaded
        embeddings = model.encode(chunks, show_progress_bar=False)
        if isinstance(embeddings, np.ndarray):
            if embeddings.ndim == 1:
                embeddings = embeddings.reshape(1, -1)
            return embeddings.tolist()
        raise TypeError("[Processor] Model output was not a NumPy array.")
    except Exception as e:
        logger.error(f"[Processor] Failed to generate embeddings: {repr(e)}")
        return []

# --------------------------------------------------------------------------
# --- PUBLIC MASTER FUNCTION
# --------------------------------------------------------------------------

def process_file(file_path: Union[str, Path]) -> Dict:
    """
    Processes a single file from path to embeddings.

    This function orchestrates the extraction, cleaning, chunking, and embedding
    of a file's content.

    Args:
        file_path: The path to the file to be processed.

    Returns:
        A dictionary containing all processed data, including embeddings.
    """
    path = Path(file_path)
    if not path.is_file():
        logger.warning(f"[Processor] Invalid file path: {file_path}")
        raise FileNotFoundError(f"Path does not point to a file: {file_path}")

    # 1. Extract
    file_type = path.suffix.lower().lstrip('.')
    file_name = path.stem
    parent_folder_path = path.parent.resolve()
    parent_folder_name = path.parent.name
    logger.debug(f"[Processor] Extracting content from: {path.name}")
    raw_content = _extract_content(path, file_type)
    content_hash = _compute_hash(raw_content)

    # 2. Clean
    cleaned_content = _clean_text(raw_content)
    cleaned_file_name = _clean_text(file_name)

    # 3. Chunk
    chunks = _chunk_text(cleaned_content)
    if not chunks:
        logger.warning(f"[Processor] No content to embed for: {path.name}")

    # 4. Embed
    embeddings = _embed_texts(chunks)
    if not embeddings:
        logger.warning(f"[Processor] Embedding failed or returned empty for: {path.name}")

    logger.info(f"[Processor] {path.name} | Type: {file_type} | Hash: {content_hash[:8]} | Chunks: {len(chunks)}")

    return {
        "file_path": str(path),
        "content": cleaned_content, # Cleaned content for potential future use
        "file_name": cleaned_file_name,
        "parent_folder": parent_folder_name, # Original parent folder
        "parent_folder_path": str(parent_folder_path),
        "file_type": file_type,
        "content_hash": content_hash,
        "embeddings": embeddings,
    }